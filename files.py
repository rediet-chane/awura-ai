import io

file_store: dict[str, str] = {}

def extract_text(filename: str, content: bytes) -> str:
    name = filename.lower()
    try:
        if name.endswith(('.txt', '.md')):
            return content.decode("utf-8", errors="ignore")

        elif name.endswith('.pdf'):
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(content))
            return "\n".join(p.extract_text() or "" for p in reader.pages)

        elif name.endswith('.docx'):
            import docx
            doc = docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)

        elif name.endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            lines.append(para.text)
            return "\n".join(lines)

        else:
            return content.decode("utf-8", errors="ignore")

    except Exception as e:
        return f"[Error reading {filename}: {e}]"

def save_file(filename: str, content: bytes) -> int:
    text = extract_text(filename, content)
    file_store[filename] = text
    return len(text)

def get_file_context(filenames: list[str]) -> str:
    ctx = ""
    for name in filenames:
        if name in file_store:
            ctx += f"\n--- {name} ---\n{file_store[name]}\n"
    return ctx

def remove_file(filename: str):
    file_store.pop(filename, None)

def list_files() -> list[str]:
    return list(file_store.keys())